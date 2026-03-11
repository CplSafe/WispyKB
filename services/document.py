# DocumentProcessor - 文档解析与分块处理服务
# 从 main_pgvector.py 拆分

import csv
import hashlib
import json
import logging
import re
import xml.etree.ElementTree as ET
from pathlib import Path
from typing import Dict, List, Any, Optional

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """文档解析与分块处理器 - 支持多种文件格式"""

    # 支持的文件扩展名
    SUPPORTED_EXTENSIONS = {
        # 文本文件
        '.txt', '.md', '.markdown', '.rst', '.log',
        # 结构化数据
        '.json', '.xml', '.yaml', '.yml', '.csv', '.toml', '.ini',
        # Office 文档
        '.docx', '.xlsx', '.pptx',
        # 网页文件
        '.html', '.htm', '.xhtml',
        # PDF (需要额外配置)
        '.pdf',
    }

    # MIME 类型映射
    MIME_TYPE_MAP = {
        'text/plain': '.txt',
        'text/markdown': '.md',
        'application/json': '.json',
        'application/xml': '.xml',
        'text/xml': '.xml',
        'text/csv': '.csv',
        'text/html': '.html',
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document': '.docx',
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': '.xlsx',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation': '.pptx',
        'application/pdf': '.pdf',
    }

    def __init__(self, html2text_available=None, bs4_available=None,
                 docx_available=None, pptx_available=None, xlsx_available=None):
        # 若未显式传入，则自动探测
        def _check(flag, module_name):
            if flag is not None:
                return flag
            try:
                __import__(module_name)
                return True
            except ImportError:
                return False

        self.html2text_available = _check(html2text_available, 'html2text')
        self.bs4_available       = _check(bs4_available,       'bs4')
        self.docx_available      = _check(docx_available,      'docx')
        self.pptx_available      = _check(pptx_available,      'pptx')
        self.xlsx_available      = _check(xlsx_available,      'openpyxl')

        self.html_converter = None
        if self.html2text_available:
            try:
                import html2text
                self.html_converter = html2text.HTML2Text()
                self.html_converter.ignore_links = False
                self.html_converter.ignore_images = False
                self.html_converter.body_width = 0
            except:
                pass

    async def process(
        self,
        file_path: str,
        filename: str,
        kb_id: str,
        doc_id: str,
        pool_ref,
        embedding_service,
        chunk_size: int = 512,
        chunk_overlap: int = 50,
        incremental: bool = False,
        vector_store_type: str = 'pgvector',
        vector_store_instance=None,
        on_progress=None
    ) -> Dict[str, Any]:
        """处理文档文件

        Args:
            on_progress: 可选进度回调 async(progress: float, message: str) -> None
                         progress 范围 0-100
        """

        async def _report(progress: float, message: str):
            if on_progress:
                try:
                    await on_progress(progress, message)
                except Exception:
                    pass

        # 获取文件扩展名
        ext = Path(filename).suffix.lower()

        # 根据文件类型读取内容
        await _report(5, "正在读取文件...")
        try:
            raw_content = await self._read_file(file_path, ext, doc_id)
        except Exception as e:
            return {"status": "failed", "error": f"读取文件失败: {str(e)}"}

        if not raw_content:
            return {"status": "failed", "error": "文件内容为空"}

        # 计算文件哈希（用于增量更新检测）
        content_bytes = raw_content.encode('utf-8') if isinstance(raw_content, str) else raw_content
        file_hash = hashlib.md5(content_bytes).hexdigest()

        # 检查是否需要增量更新
        if incremental:
            async with pool_ref.connection() as conn:
                from psycopg.rows import dict_row
                async with conn.cursor(row_factory=dict_row) as cur:
                    await cur.execute("SELECT file_hash FROM documents WHERE id = %s", (doc_id,))
                    row = await cur.fetchone()
                    if row and row['file_hash'] == file_hash:
                        # 文档未变更，但确保状态是completed
                        await cur.execute("UPDATE documents SET status = 'completed' WHERE id = %s", (doc_id,))
                        await conn.commit()
                        return {"status": "unchanged", "message": "文档未变更，无需更新"}

        # 解析内容
        await _report(10, "正在解析文档内容...")
        content = self._parse_content(raw_content, ext)

        # 智能分块
        await _report(15, "正在进行智能分块...")
        chunks = self._chunk_content(content, chunk_size, chunk_overlap)

        await _report(20, f"分块完成，共 {len(chunks)} 个分块，正在保存...")

        # 处理数据库事务 - 第一阶段：保存分块到数据库
        chunk_ids = []
        chunk_contents = []  # 保存 chunk 内容用于后续 Milvus 存储

        async with pool_ref.connection() as conn:
            async with conn.cursor() as cur:
                # 清除旧分块（增量更新）
                await cur.execute("DELETE FROM chunks WHERE doc_id = %s", (doc_id,))

                # 插入新分块
                for idx, chunk_content in enumerate(chunks):
                    chunk_id = f"{doc_id}-{idx}"
                    await cur.execute("""
                        INSERT INTO chunks (id, doc_id, chunk_index, content, metadata)
                        VALUES (%s, %s, %s, %s, %s)
                        RETURNING id
                    """, (chunk_id, doc_id, idx, chunk_content, json.dumps({"size": len(chunk_content)})))
                    chunk_ids.append(chunk_id)
                    chunk_contents.append(chunk_content)  # 保存内容

                # 更新文档
                await cur.execute("""
                    UPDATE documents
                    SET status = 'completed', chunk_count = %s, content = %s, file_hash = %s, updated_at = NOW()
                    WHERE id = %s
                """, (len(chunks), raw_content[:10000], file_hash, doc_id))  # 只存储前10000字符预览

            await conn.commit()

        # 生成向量嵌入 - 逐个生成并报告细粒度进度
        # embedding 阶段占总进度的 25% ~ 85%（共 60%）
        await _report(25, f"开始生成向量嵌入 (0/{len(chunks)})...")
        embeddings = []
        total_chunks = len(chunks)
        for i, text in enumerate(chunks):
            embedding = await embedding_service.generate(text)
            embeddings.append(embedding)
            # 25% ~ 85% 区间内线性分布
            embed_progress = 25 + (i + 1) / total_chunks * 60
            await _report(embed_progress, f"正在生成向量嵌入 ({i + 1}/{total_chunks})...")

        # 根据向量存储类型存储向量
        await _report(87, "正在存储向量数据...")
        valid_embeddings = [(chunk_id, emb) for chunk_id, emb in zip(chunk_ids, embeddings) if emb]

        if vector_store_type == 'milvus' and vector_store_instance:
            # 使用 Milvus 存储向量
            if valid_embeddings:
                items = []
                for idx, (chunk_id, embedding) in enumerate(valid_embeddings):
                    # 直接使用之前保存的内容，不需要再次查询数据库
                    items.append({
                        'chunk_id': chunk_id,
                        'kb_id': kb_id,
                        'document_id': doc_id,
                        'content': chunk_contents[idx],
                        'embedding': embedding,
                        'chunk_index': int(chunk_id.split('-')[-1]),
                    })

                if items:
                    count = await vector_store_instance.insert_batch(items)
                    logger.info(f"向量嵌入存储到 Milvus 完成: {count}/{len(chunks)} 个分块")
            else:
                logger.warning(f"向量嵌入生成失败，文档将保存但不包含向量数据")

        else:
            # 使用 pgvector 存储向量（PostgreSQL）
            async with pool_ref.connection() as conn:
                async with conn.cursor() as cur:
                    if valid_embeddings:
                        for chunk_id, embedding in valid_embeddings:
                            await cur.execute(
                                "UPDATE chunks SET embedding = %s::vector WHERE id = %s",
                                (embedding, chunk_id)
                            )
                        await conn.commit()
                        logger.info(f"向量嵌入生成完成: {len(valid_embeddings)}/{len(chunks)} 个分块")
                    else:
                        logger.warning(f"向量嵌入生成失败，文档将保存但不包含向量数据")

        await _report(100, f"文档处理完成，共 {len(chunks)} 个分块")
        return {
            "status": "success",
            "doc_id": doc_id,
            "chunk_count": len(chunks),
            "message": f"文档处理完成，共 {len(chunks)} 个分块"
        }

    async def _read_file(self, file_path: str, ext: str, doc_id: str = None) -> str:
        """根据文件类型读取文件内容"""
        file_path_obj = Path(file_path)

        # 文本文件直接读取
        if ext in {'.txt', '.md', '.markdown', '.rst', '.log', '.yaml', '.yml', '.toml', '.ini'}:
            encodings = ['utf-8', 'gbk', 'gb2312', 'latin-1']
            for enc in encodings:
                try:
                    with open(file_path_obj, 'r', encoding=enc) as f:
                        return f.read()
                except UnicodeDecodeError:
                    continue
            raise ValueError("无法解码文件，尝试的编码: " + ", ".join(encodings))

        # JSON 文件
        elif ext == '.json':
            with open(file_path_obj, 'r', encoding='utf-8') as f:
                data = json.load(f)
                return json.dumps(data, ensure_ascii=False, indent=2)

        # CSV 文件
        elif ext == '.csv':
            with open(file_path_obj, 'r', encoding='utf-8') as f:
                reader = csv.DictReader(f)
                rows = list(reader)
                return "\n".join([", ".join(row.values()) for row in rows])

        # XML 文件
        elif ext == '.xml':
            with open(file_path_obj, 'r', encoding='utf-8') as f:
                content = f.read()
            # 尝试美化 XML
            try:
                root = ET.fromstring(content)
                return ET.tostring(root, encoding='unicode', method='xml')
            except:
                return content

        # HTML 文件
        elif ext in {'.html', '.htm', '.xhtml'}:
            if self.html_converter:
                with open(file_path_obj, 'r', encoding='utf-8') as f:
                    html_content = f.read()
                return self.html_converter.handle(html_content)
            elif self.bs4_available:
                from bs4 import BeautifulSoup
                with open(file_path_obj, 'r', encoding='utf-8') as f:
                    html_content = f.read()
                soup = BeautifulSoup(html_content, 'html.parser')
                # 移除脚本和样式
                for script in soup(["script", "style"]):
                    script.decompose()
                return soup.get_text(separator="\n", strip=True)
            else:
                with open(file_path_obj, 'r', encoding='utf-8') as f:
                    return f.read()

        # DOCX 文件
        elif ext == '.docx':
            if not self.docx_available:
                raise ValueError("DOCX 支持未安装，请运行: pip install python-docx")
            from docx import Document as DocxDocument
            doc = DocxDocument(file_path_obj)
            paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text.strip())
            # 提取表格
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells])
                    if row_text.strip():
                        paragraphs.append(row_text)
            return "\n\n".join(paragraphs)

        # XLSX 文件 - 支持图片提取（包括 WPS 格式）
        elif ext == '.xlsx':
            if not self.xlsx_available:
                raise ValueError("XLSX 支持未安装，请运行: pip install openpyxl")
            import openpyxl
            from core import config
            import zipfile
            import xml.etree.ElementTree as ET

            file_path_str = str(file_path_obj)

            # 初始化图片映射: {image_id: image_url}
            images_map = {}
            images_dir = None

            # 获取静态文件基础 URL（完整 URL）
            static_base_url = config.STATIC_URL

            if doc_id:
                # 创建图片存储目录
                images_dir = Path(config.UPLOAD_DIR) / "images" / doc_id
                images_dir.mkdir(parents=True, exist_ok=True)

                # 1. 提取图片并建立映射（处理 WPS 格式）
                try:
                    with zipfile.ZipFile(file_path_str, 'r') as zip_ref:
                        # 读取 cellimages.xml.rels 获取 rId 到图片文件的映射
                        rid_to_image = {}
                        try:
                            rels_content = zip_ref.read('xl/_rels/cellimages.xml.rels').decode('utf-8')
                            rels_root = ET.fromstring(rels_content)
                            ns = {'r': 'http://schemas.openxmlformats.org/package/2006/relationships'}
                            for rel in rels_root.findall('r:Relationship', ns):
                                rid = rel.get('Id')
                                target = rel.get('Target')
                                rid_to_image[rid] = target
                        except Exception as e:
                            logger.warning(f"读取 cellimages.xml.rels 失败: {e}")

                        # 读取 cellimages.xml 获取图片 ID 到 rId 的映射
                        try:
                            cellimages_content = zip_ref.read('xl/cellimages.xml').decode('utf-8')
                            cellimages_root = ET.fromstring(cellimages_content)
                            ns = {'etc': 'http://www.wps.cn/officeDocument/2017/etCustomData',
                                  'xdr': 'http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing',
                                  'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
                                  'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}

                            for cell_image in cellimages_root.findall('etc:cellImage', ns):
                                pic = cell_image.find('xdr:pic', ns)
                                if pic is not None:
                                    nv_pic_pr = pic.find('xdr:nvPicPr', ns)
                                    if nv_pic_pr is not None:
                                        c_nv_pr = nv_pic_pr.find('xdr:cNvPr', ns)
                                        if c_nv_pr is not None:
                                            image_id = c_nv_pr.get('name')  # ID_xxx 格式

                                            # 找到 r:embed
                                            blip_fill = pic.find('xdr:blipFill', ns)
                                            if blip_fill is not None:
                                                blip = blip_fill.find('a:blip', ns)
                                                if blip is not None:
                                                    rid = blip.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed')

                                                    # 获取图片文件名
                                                    if rid in rid_to_image:
                                                        image_filename = rid_to_image[rid]

                                                        # 提取并保存图片
                                                        try:
                                                            image_data = zip_ref.read(f'xl/{image_filename}')
                                                            saved_filename = f"{Path(image_filename).stem}.png"
                                                            saved_path = images_dir / saved_filename

                                                            with open(saved_path, 'wb') as f:
                                                                f.write(image_data)

                                                            # 使用完整的 HTTP URL
                                                            image_url = f"{static_base_url}/images/{doc_id}/{saved_filename}"
                                                            images_map[image_id] = image_url
                                                            logger.info(f"保存图片: {image_id} -> {image_url}")
                                                        except Exception as e:
                                                            logger.warning(f"保存图片失败 {image_id}: {e}")
                        except Exception as e:
                            logger.warning(f"读取 cellimages.xml 失败: {e}")
                except Exception as e:
                    logger.warning(f"提取图片失败: {e}")

            # 2. 读取工作表内容（保留公式以获取 DISPIMG）
            wb = openpyxl.load_workbook(file_path_obj, data_only=False)
            sheets_data = []

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                sheet_data = [f"## 工作表: {sheet_name}"]

                # 遍历每一行
                for row_idx, row in enumerate(ws.iter_rows(values_only=False)):
                    row_data = []
                    for col_idx, cell in enumerate(row):
                        # 获取单元格值
                        if cell.data_type == 'f':  # 公式
                            cell_value = cell.value
                            # 处理 DISPIMG 公式
                            if cell_value and 'DISPIMG' in str(cell_value):
                                # 提取图片 ID: =_xlfn.DISPIMG("ID_xxx",1)
                                import re
                                match = re.search(r'DISPIMG\("([^"]+)"', str(cell_value))
                                if match:
                                    image_id = match.group(1)
                                    if image_id in images_map:
                                        cell_value = images_map[image_id]
                                    else:
                                        cell_value = "[图片]"
                                else:
                                    cell_value = "[图片]"
                            else:
                                # 其他公式，计算值
                                try:
                                    cell_value = str(cell.value) if cell.value else ""
                                except:
                                    cell_value = ""
                        else:
                            cell_value = str(cell.value) if cell.value is not None else ""

                        row_data.append(cell_value)

                    if any(row_data):  # 跳过空行
                        # 检查是否有图片 URL，如果有则使用 Markdown 表格格式
                        has_image = any('/static/files/images/' in str(cell) for cell in row_data)
                        if has_image:
                            # 将图片 URL 转换为 Markdown 图片格式
                            formatted_cells = []
                            for cell in row_data:
                                cell_str = str(cell)
                                if '/static/files/images/' in cell_str:
                                    cell_str = f"![图片]({cell_str})"
                                formatted_cells.append(cell_str)
                            sheet_data.append("| " + " | ".join(formatted_cells) + " |")
                        else:
                            sheet_data.append(" | ".join(row_data))

                sheets_data.append("\n".join(sheet_data))

            wb.close()
            return "\n\n".join(sheets_data)

        # PPTX 文件
        elif ext == '.pptx':
            if not self.pptx_available:
                raise ValueError("PPTX 支持未安装，请运行: pip install python-pptx")
            from pptx import Presentation
            prs = Presentation(file_path_obj)
            slides_text = []
            for i, slide in enumerate(prs.slides):
                slide_text = [f"## 幻灯片 {i + 1}"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    # 处理表格
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_data = [cell.text.strip() for cell in row.cells]
                            slide_text.append(" | ".join(row_data))
                if len(slide_text) > 1:
                    slides_text.append("\n".join(slide_text))
            return "\n\n".join(slides_text)

        # PDF 文件 (高级解析 - 使用 PyMuPDF)
        elif ext == '.pdf':
            # 使用增强的 PDF 解析器
            try:
                from .pdf_parser import PDFParser
                parser = PDFParser()
                if not parser.available:
                    # 回退到基础解析
                    raise ImportError("PyMuPDF 未安装")
                # 返回解析结果供后续处理
                return f"__PDF_PARSED__:{file_path}"
            except Exception as e:
                logger.warning(f"高级 PDF 解析失败，使用基础解析: {e}")
                # 回退到基础解析
                try:
                    import PyPDF2
                    with open(file_path_obj, 'rb') as f:
                        reader = PyPDF2.PdfReader(f)
                        text_parts = []
                        for page in reader.pages:
                            text = page.extract_text()
                            if text:
                                text_parts.append(text)
                        return "\n\n".join(text_parts)
                except ImportError:
                    raise ValueError("PDF 支持未安装，请运行: pip install PyMuPDF 或 PyPDF2")
                except Exception as e2:
                    raise ValueError(f"PDF 解析失败: {str(e2)}")

        else:
            # 尝试作为文本读取
            try:
                with open(file_path_obj, 'r', encoding='utf-8') as f:
                    return f.read()
            except:
                raise ValueError(f"不支持的文件类型: {ext}")

    def _parse_content(self, content: str, ext: str) -> str:
        """根据文件类型解析和清理内容"""
        # 清理多余空白
        lines = content.split('\n')
        cleaned_lines = []
        for line in lines:
            stripped = line.strip()
            if stripped:
                cleaned_lines.append(stripped)

        parsed_content = "\n\n".join(cleaned_lines)

        # 对于 Excel 文件，额外处理 DISPIMG 占位符
        if ext == '.xlsx':
            parsed_content = self._clean_dispimg_placeholders(parsed_content)

        return parsed_content

    def _clean_dispimg_placeholders(self, content: str) -> str:
        """
        清理 Excel 中的 DISPIMG 公式占位符

        将 =DISPIMG("ID_xxx",1) 替换为友好的图片标记
        同时保留上下文信息以便理解

        Args:
            content: 原始内容

        Returns:
            清理后的内容
        """
        import re

        # 匹配 DISPIMG 公式: =DISPIMG("ID_xxx", N)
        def replace_dispimg(match):
            # 提取图片ID
            img_id = match.group(1)
            # 生成简短标识
            short_id = img_id[-8:] if len(img_id) > 8 else img_id
            return f"[图片_{short_id}]"

        # 替换所有 DISPIMG 公式
        pattern = r'=DISPIMG\("([^"]+)",\s*\d+\)'
        content = re.sub(pattern, replace_dispimg, content)

        return content

    def _chunk_content(self, content: str, chunk_size: int, chunk_overlap: int) -> List[str]:
        """
        智能分块 - 优先使用智能分块器

        支持：
        1. Markdown 标题识别 (H1-H6)
        2. 父链保留（每个 chunk 保留父级标题上下文）
        3. 关键词提取
        4. 智能断句
        """
        # 检查是否是 PDF 解析后的结果
        if content.startswith("__PDF_PARSED__:"):
            file_path = content.replace("__PDF_PARSED__:", "")
            from .pdf_parser import PDFParser
            parser = PDFParser()
            result = parser.parse(file_path, limit=chunk_size)
            # 返回 {"title": "标题", "content": "内容"} 格式
            # 需要转换为字符串列表
            return [f"{c.get('title', '')}\n{c['content']}" for c in result.get("content", [])]

        # 尝试使用智能分块器
        try:
            from .smart_chunk import SmartChunker
            chunker = SmartChunker(limit=chunk_size, extract_keywords=True)
            chunks_data = chunker.parse(content)

            # 转换为字符串列表（保留标题信息）
            chunks = []
            for c in chunks_data:
                title = c.get("title", "")
                body = c.get("content", "")
                keywords = c.get("keywords", [])

                # 组合标题和内容
                if title:
                    chunk = f"# {title}\n\n{body}"
                else:
                    chunk = body

                # 如果有关键词，可以添加到元数据
                if keywords:
                    # 这里可以添加关键词到元数据，暂时不添加到内容中
                    pass

                chunks.append(chunk)

            if chunks:
                return chunks
        except Exception as e:
            logger.warning(f"智能分块失败，使用基础分块: {e}")

        # 回退到基础分块
        return self._basic_chunk(content, chunk_size, chunk_overlap)

    def _basic_chunk(self, content: str, chunk_size: int, chunk_overlap: int) -> List[str]:
        """基础分块 - 基于段落"""
        chunks = []

        # 按段落分割
        paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]

        current_chunk = ""
        current_size = 0

        for para in paragraphs:
            para_size = len(para)

            if current_size + para_size > chunk_size and current_chunk:
                chunks.append(current_chunk.strip())
                # 保留重叠部分
                overlap_text = current_chunk[-chunk_overlap:] if len(current_chunk) > chunk_overlap else current_chunk
                current_chunk = overlap_text + "\n\n" + para
                current_size = len(current_chunk)
            else:
                if current_chunk:
                    current_chunk += "\n\n" + para
                else:
                    current_chunk = para
                current_size = len(current_chunk)

        if current_chunk:
            chunks.append(current_chunk.strip())

        # 如果没有段落，按字符分割
        if not chunks:
            for i in range(0, len(content), chunk_size - chunk_overlap):
                chunks.append(content[i:i + chunk_size])

        return chunks
